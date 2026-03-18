"""
Provision PPTX Dubbing - Translate PPTX
========================================
Extracts text from a PowerPoint presentation, translates it using Gemini AI,
and produces a translated PPTX file.

Usage:
    python translate_pptx.py --config path/to/config.json
"""

import argparse
import json
import os
import copy
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt

# TODO: Install google-generativeai: pip install google-generativeai
import google.generativeai as genai


def load_config(config_path: str) -> dict:
    """Load project configuration from JSON file."""
    with open(config_path, "r", encoding="utf-8") as f:
        return json.load(f)


def resolve_path(project_dir: str, relative_path: str) -> str:
    """Resolve a path relative to the project directory."""
    return os.path.join(project_dir, relative_path)


def extract_slide_texts(prs: Presentation) -> list[dict]:
    """
    Extract all text from each slide in the presentation.

    Returns a list of dicts, one per slide:
    [
        {
            "slide_index": 0,
            "texts": [
                {"shape_name": "Title 1", "text": "Welcome to Provision ISR", "runs": [...]},
                ...
            ]
        },
        ...
    ]
    """
    slides_data = []

    for slide_idx, slide in enumerate(prs.slides):
        slide_texts = []

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            shape_text_parts = []
            for paragraph in shape.text_frame.paragraphs:
                runs_data = []
                for run in paragraph.runs:
                    if run.text.strip():
                        runs_data.append({
                            "text": run.text,
                            "font_name": run.font.name,
                            "font_size": run.font.size,
                            "bold": run.font.bold,
                            "italic": run.font.italic,
                        })
                if runs_data:
                    shape_text_parts.append({
                        "paragraph_runs": runs_data,
                        "full_text": " ".join(r["text"] for r in runs_data),
                    })

            if shape_text_parts:
                full_shape_text = "\n".join(p["full_text"] for p in shape_text_parts)
                slide_texts.append({
                    "shape_name": shape.name,
                    "shape_id": shape.shape_id,
                    "text": full_shape_text,
                    "paragraphs": shape_text_parts,
                })

        slides_data.append({
            "slide_index": slide_idx,
            "texts": slide_texts,
        })

    return slides_data


def translate_texts_with_gemini(
    slides_data: list[dict],
    target_language: str,
    api_key: str,
    model_name: str = "gemini-2.0-flash",
) -> list[dict]:
    """
    Translate all extracted slide texts using Gemini AI.

    Sends the full text content to Gemini with instructions to preserve
    formatting structure and return translations in the same JSON structure.
    """
    genai.configure(api_key=api_key)
    model = genai.GenerativeModel(model_name)

    # Build a structured payload for Gemini
    texts_to_translate = []
    for slide in slides_data:
        for shape in slide["texts"]:
            for para in shape["paragraphs"]:
                for run in para["paragraph_runs"]:
                    texts_to_translate.append(run["text"])

    if not texts_to_translate:
        print("No text found to translate.")
        return slides_data

    # Batch translate in chunks to avoid token limits
    chunk_size = 100
    translated_texts = []

    for i in range(0, len(texts_to_translate), chunk_size):
        chunk = texts_to_translate[i : i + chunk_size]

        prompt = f"""Translate the following texts from English to {target_language}.
Return a JSON array of translated strings in the EXACT same order.
Preserve any technical terms, product names (like "Provision ISR", "NVR", "DVR"),
and model numbers as-is. Do not add explanations.

Texts to translate:
{json.dumps(chunk, ensure_ascii=False, indent=2)}

Return ONLY a JSON array of translated strings, nothing else."""

        response = model.generate_content(prompt)
        response_text = response.text.strip()

        # Parse the JSON array from the response
        # Strip markdown code fences if present
        if response_text.startswith("```"):
            lines = response_text.split("\n")
            response_text = "\n".join(lines[1:-1])

        try:
            chunk_translated = json.loads(response_text)
            translated_texts.extend(chunk_translated)
        except json.JSONDecodeError:
            print(f"WARNING: Failed to parse Gemini response for chunk {i}. Using originals.")
            translated_texts.extend(chunk)

    # Map translated texts back into the slide data structure
    translated_data = json.loads(json.dumps(slides_data))  # deep copy
    text_idx = 0

    for slide in translated_data:
        for shape in slide["texts"]:
            for para in shape["paragraphs"]:
                for run in para["paragraph_runs"]:
                    if text_idx < len(translated_texts):
                        run["translated_text"] = translated_texts[text_idx]
                        text_idx += 1

    return translated_data


def apply_translations_to_pptx(prs: Presentation, translated_data: list[dict]) -> Presentation:
    """
    Replace text in the presentation with translated text while preserving formatting.
    """
    for slide_data in translated_data:
        slide_idx = slide_data["slide_index"]
        slide = prs.slides[slide_idx]

        # Build a lookup by shape_id
        shape_lookup = {s["shape_id"]: s for s in slide_data["texts"]}

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if shape.shape_id not in shape_lookup:
                continue

            shape_data = shape_lookup[shape.shape_id]
            para_idx = 0

            for paragraph in shape.text_frame.paragraphs:
                run_idx = 0
                for run in paragraph.runs:
                    if not run.text.strip():
                        continue

                    # Find the matching translated text
                    if para_idx < len(shape_data["paragraphs"]):
                        para_data = shape_data["paragraphs"][para_idx]
                        if run_idx < len(para_data["paragraph_runs"]):
                            run_data = para_data["paragraph_runs"][run_idx]
                            if "translated_text" in run_data:
                                run.text = run_data["translated_text"]
                            run_idx += 1

                # Move to next paragraph data if we consumed all runs
                if run_idx > 0:
                    para_idx += 1

    return prs


def save_narration_texts(translated_data: list[dict], output_path: str):
    """
    Save the translated text per slide for narration generation.
    Each slide gets a single combined text string that will be spoken.
    """
    narration_texts = []

    for slide_data in translated_data:
        slide_text_parts = []
        for shape in slide_data["texts"]:
            for para in shape["paragraphs"]:
                for run in para["paragraph_runs"]:
                    text = run.get("translated_text", run["text"])
                    slide_text_parts.append(text)

        narration_texts.append({
            "slide_index": slide_data["slide_index"],
            "narration_text": " ".join(slide_text_parts),
        })

    with open(output_path, "w", encoding="utf-8") as f:
        json.dump(narration_texts, f, ensure_ascii=False, indent=2)

    print(f"Saved narration texts to {output_path}")


def main():
    parser = argparse.ArgumentParser(description="Translate PPTX presentation")
    parser.add_argument("--config", required=True, help="Path to config.json")
    args = parser.parse_args()

    config = load_config(args.config)
    project_dir = config["project_dir"]

    # Resolve paths
    source_path = resolve_path(project_dir, config["source_pptx"])
    translated_path = resolve_path(project_dir, config["translated_pptx"])
    narration_texts_path = resolve_path(project_dir, config["narration_dir"] + "/narration_texts.json")

    # Ensure output directories exist
    os.makedirs(os.path.dirname(translated_path), exist_ok=True)
    os.makedirs(os.path.dirname(narration_texts_path), exist_ok=True)

    print(f"Loading presentation: {source_path}")
    prs = Presentation(source_path)

    print("Extracting text from slides...")
    slides_data = extract_slide_texts(prs)
    total_texts = sum(len(s["texts"]) for s in slides_data)
    print(f"  Found {total_texts} text shapes across {len(slides_data)} slides")

    print(f"Translating to {config['target_language']}...")
    translated_data = translate_texts_with_gemini(
        slides_data,
        target_language=config["target_language"],
        api_key=config["gemini_api_key"],
        model_name=config.get("gemini_model_translation", "gemini-2.0-flash"),
    )

    print("Applying translations to PPTX...")
    prs = apply_translations_to_pptx(prs, translated_data)

    print(f"Saving translated PPTX: {translated_path}")
    prs.save(translated_path)

    # Save narration texts for the next pipeline stage
    save_narration_texts(translated_data, narration_texts_path)

    print("Translation complete.")


if __name__ == "__main__":
    main()
